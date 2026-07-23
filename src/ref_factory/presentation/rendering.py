"""
Rendu PPTX one-slide conforme à la charte Orange Cyberdefense.

La fiche est TOUJOURS générée en remplissant le MÊME template exact — une
fiche REF réelle validée, embarquée dans le repo :
    src/ref_factory/templates/fiche_ref_template.pptx
Quel que soit le secteur ou le contenu, le design ne change jamais : photo,
mise en page, couleurs, polices et bloc "Profil de la prestation" restent
identiques ; seuls les textes sont remplacés par les informations données
par l'utilisateur (le logo du client d'origine est retiré et remplacé par
le nom du client). Les fiches du RAG ne servent qu'à inspirer le contenu
texte, jamais la forme.

Si ce template est absent (repo incomplet), un rendu de secours reconstruit
la fiche sur le template Brand Box fond noir (OFR_template_Fond_noir.potx).
"""
from __future__ import annotations

import copy
import io
import json
import re
import unicodedata
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm, Emu, Pt

from ref_factory.config import OUTPUT_DIR, PROJECT_ROOT

# ── Chargement de la charte ──────────────────────────────────────────────────
_CHARTER_PATH = Path(__file__).resolve().parents[1] / "charter" / "ocd_charter.json"


def _load_charter() -> dict:
    if _CHARTER_PATH.exists():
        return json.loads(_CHARTER_PATH.read_text(encoding="utf-8"))
    return {}


_CHARTER = _load_charter()
_COLORS = _CHARTER.get("colors", {})
_FONTS = _CHARTER.get("fonts", {})
_STYLES = _CHARTER.get("pptx", {}).get("styles", {})
_CONF_COLORS = _CHARTER.get("confidentiality_colors", {})
_DARK = _CHARTER.get("dark_theme", {})
_TEMPLATE_CFG = _CHARTER.get("template", {})

PRIMARY_FONT = _FONTS.get("primary", "Source Sans Pro")


def _font(name_override: str | None = None) -> str:
    return name_override or PRIMARY_FONT


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# Palette OCD officielle — source : Orange WHT Core.xml / Orange BLK Core.xml
ORANGE = _rgb(_COLORS.get("primary_orange", "#FF7900"))
BLACK = _rgb(_COLORS.get("black", "#000000"))
DARK_GREY = _rgb(_COLORS.get("dark_grey", "#595959"))
MEDIUM_GREY = _rgb(_COLORS.get("medium_grey", "#8F8F8F"))
LIGHT_GREY = _rgb(_COLORS.get("light_grey", "#D6D6D6"))
WHITE = _rgb(_COLORS.get("white", "#FFFFFF"))

# Déclinaison fond noir (template OFR_template_Fond_noir.potx)
DARK_BG = _rgb(_DARK.get("background", "#000000"))
TEXT_PRIMARY = _rgb(_DARK.get("text_primary", "#FFFFFF"))
TEXT_SECONDARY = _rgb(_DARK.get("text_secondary", "#D6D6D6"))
PANEL_FILL = _rgb(_DARK.get("panel_fill", "#262626"))
PANEL_BORDER = _rgb(_DARK.get("panel_border", "#595959"))

# ── Template fiche exact (chemin principal) ─────────────────────────────────

FICHE_TEMPLATE_PATH = PROJECT_ROOT / _TEMPLATE_CFG.get(
    "path", "src/ref_factory/templates/fiche_ref_template.pptx"
)

# Identifiants des formes du template fiche (stables : le fichier est versionné).
_FICHE_SHAPES = {
    "title": 2,             # Titre principal
    "context_body": 25,     # Corps de la section Contexte
    "realisation_body": 26, # Corps de la section Réalisation (puces)
    "benefits_body": 27,    # Corps de la section Bénéfices (puces)
    "year": 66,             # Pastille année ("2024")
    "sector": 67,           # Pastille secteur ("Transport")
    "location": 68,         # Pastille lieu ("Paris / Colombie")
    "duration": 71,         # Pastille charge ("27,25 jours")
}
_FICHE_CLIENT_LOGO_ID = 1028  # Logo du client d'origine, remplacé par son nom

# ── Template Brand Box fond noir (rendu de secours) ──────────────────────────

TEMPLATE_PATH = PROJECT_ROOT / _TEMPLATE_CFG.get(
    "fallback_potx",
    "Tools and templates PPT - FR/French/2. Templates/French/OFR_template_Fond_noir.potx",
)
TEMPLATE_LAYOUT_NAME = _TEMPLATE_CFG.get("layout_name", "Vide")

# Dimensions du template fond noir (16:9) — utilisées aussi en fallback
DESIGN_W = Cm(33.87)  # grille de conception d'origine, mise à l'échelle du template
FALLBACK_W = Cm(_TEMPLATE_CFG.get("slide_width_cm", 25.40))
FALLBACK_H = Cm(_TEMPLATE_CFG.get("slide_height_cm", 14.29))


def _load_template_presentation() -> Presentation | None:
    """Charge le .potx OCD comme une présentation (patch du content-type en mémoire)."""
    if not TEMPLATE_PATH.exists():
        return None
    try:
        data = TEMPLATE_PATH.read_bytes()
        src = zipfile.ZipFile(io.BytesIO(data))
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.namelist():
                content = src.read(item)
                if item == "[Content_Types].xml":
                    content = content.replace(
                        b"presentationml.template.main+xml",
                        b"presentationml.presentation.main+xml",
                    )
                dst.writestr(item, content)
        buffer.seek(0)
        return Presentation(buffer)
    except Exception:
        return None


def _strip_example_slides(prs: Presentation) -> None:
    """Retire les slides d'exemple du template (le master et les layouts restent)."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _find_layout(prs: Presentation):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name.strip().lower() == TEMPLATE_LAYOUT_NAME.strip().lower():
                return layout
    return prs.slide_masters[0].slide_layouts[-1]


def _new_fiche_slide():
    """Retourne (présentation, slide vierge) basées sur le template OCD fond noir.

    Le template est identique pour toutes les fiches. Si le dossier Brand Box
    est absent (repo incomplet), un fallback vierge aux mêmes dimensions et au
    même fond sombre est utilisé pour ne pas bloquer la génération.
    """
    prs = _load_template_presentation()
    if prs is not None:
        _strip_example_slides(prs)
        slide = prs.slides.add_slide(_find_layout(prs))
    else:
        prs = Presentation()
        prs.slide_width = FALLBACK_W
        prs.slide_height = FALLBACK_H
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Fond noir explicite — même mécanisme que les slides du template officiel
    # (override de fond par slide ; tx1 = noir dans le thème Orange).
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    return prs, slide


# ── Grille proportionnelle ───────────────────────────────────────────────────

class _Grid:
    """Géométrie de la fiche, proportionnelle aux dimensions réelles du template.

    La grille d'origine a été conçue pour 33.87 cm de large ; le facteur
    d'échelle la transpose à l'identique sur le template fond noir (25.40 cm)
    ou tout autre format 16:9.
    """

    def __init__(self, slide_w: int, slide_h: int) -> None:
        self.w = slide_w
        self.h = slide_h
        self.scale = slide_w / DESIGN_W

    def cm(self, value: float) -> Emu:
        return Emu(int(Cm(value) * self.scale))

    def pt(self, value: float) -> Pt:
        return Pt(max(6, round(value * self.scale)))

    @property
    def margin(self) -> Emu:
        return self.cm(1.2)

    @property
    def bandeau_h(self) -> Emu:
        return self.cm(0.6)

    @property
    def left_panel_w(self) -> Emu:
        return self.cm(9.0)

    @property
    def right_x(self) -> Emu:
        return Emu(self.margin + self.left_panel_w + self.cm(0.5))

    @property
    def right_w(self) -> Emu:
        return Emu(self.w - self.right_x - self.margin)

    @property
    def content_top(self) -> Emu:
        return Emu(self.bandeau_h + self.cm(3.5) + self.cm(0.3))

    @property
    def content_h(self) -> Emu:
        return Emu(self.h - self.content_top - self.cm(1.5))


# ── Rendu principal ──────────────────────────────────────────────────────────

def render_ref_slide(payload: dict[str, Any]) -> str:
    """Génère une slide REF unique en remplissant le template fiche exact.

    Le template (photo, design, bloc profil, polices, couleurs) est toujours
    le même ; seuls les textes issus du payload changent. Rendu de secours
    programmatique sur le template fond noir si le fichier est absent.
    """
    if FICHE_TEMPLATE_PATH.exists():
        return _render_from_fiche_template(payload)
    return _render_programmatic(payload)


# ── Remplissage du template fiche exact ──────────────────────────────────────

def _walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # groupe
            yield from _walk_shapes(shape.shapes)


def _render_from_fiche_template(payload: dict[str, Any]) -> str:
    prs = Presentation(str(FICHE_TEMPLATE_PATH))
    slide = prs.slides[0]
    shapes = {sh.shape_id: sh for sh in _walk_shapes(slide.shapes)}

    def fill(key: str, lines: list[str]) -> None:
        shape = shapes.get(_FICHE_SHAPES[key])
        if shape is not None and getattr(shape, "has_text_frame", False):
            _fill_text_frame(shape.text_frame, lines)

    fill("title", [_coerce(payload.get("title"), "Fiche Référence")])
    fill("context_body", _split_lines(payload.get("context")))

    realisation = _split_lines(payload.get("mission"))
    deliverables = [d for d in (payload.get("deliverables") or []) if str(d).strip()]
    realisation.extend(str(d).strip() for d in deliverables[:5])
    fill("realisation_body", realisation)

    results = [str(r).strip() for r in (payload.get("results") or []) if str(r).strip()]
    fill("benefits_body", results)

    fill("year", [_extract_year(payload)])
    fill("sector", [_coerce(payload.get("sector"), "[À COMPLÉTER]")])
    fill("location", [_coerce(payload.get("location"), "[À COMPLÉTER]")])
    fill("duration", [_coerce(payload.get("duration"), "[À COMPLÉTER]")])

    _replace_client_logo(slide, shapes, payload)

    output_path = _build_output_path(payload.get("title"))
    prs.save(output_path)
    return str(output_path)


def _fill_text_frame(tf, lines: list[str]) -> None:
    """Remplace le texte d'un cadre en préservant la mise en forme existante.

    Chaque ligne devient un paragraphe ; le style (police, taille, couleur,
    puces, alignement) est hérité des paragraphes du template — les
    paragraphes supplémentaires sont clonés du premier, les excédentaires
    supprimés.
    """
    lines = [str(line).strip() for line in (lines or []) if str(line).strip()]
    if not lines:
        lines = ["[À COMPLÉTER]"]

    paragraphs = list(tf.paragraphs)
    prototype = next((p for p in paragraphs if p.runs), paragraphs[0])
    prototype_xml = copy.deepcopy(prototype._p)
    txBody = tf._txBody

    for index, line in enumerate(lines):
        current = list(tf.paragraphs)
        if index < len(current):
            paragraph = current[index]
        else:
            txBody.append(copy.deepcopy(prototype_xml))
            paragraph = list(tf.paragraphs)[-1]
        _set_paragraph_text(paragraph, line, prototype_xml)

    for paragraph in list(tf.paragraphs)[len(lines):]:
        paragraph._p.getparent().remove(paragraph._p)


def _set_paragraph_text(paragraph, text: str, prototype_xml) -> None:
    runs = paragraph.runs
    if not runs:
        prototype_runs = prototype_xml.findall(qn("a:r"))
        if not prototype_runs:
            paragraph.text = text
            return
        paragraph._p.append(copy.deepcopy(prototype_runs[0]))
        runs = paragraph.runs
    runs[0].text = text
    for run in runs[1:]:
        run._r.getparent().remove(run._r)


def _replace_client_logo(slide, shapes: dict, payload: dict[str, Any]) -> None:
    """Retire le logo du client d'origine et affiche le nom du client à la place."""
    logo = shapes.get(_FICHE_CLIENT_LOGO_ID)
    if logo is None:
        return
    left, top, width, height = logo.left, logo.top, logo.width, logo.height
    logo._element.getparent().remove(logo._element)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = _coerce(payload.get("client"), "[CLIENT]")
    p.font.name = "Arial"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE


def _split_lines(value: Any) -> list[str]:
    if isinstance(value, list):
        return [str(v).strip() for v in value if str(v).strip()]
    if isinstance(value, str) and value.strip():
        return [part.strip() for part in value.split("\n") if part.strip()]
    return []


def _extract_year(payload: dict[str, Any]) -> str:
    for key in ("date", "year", "annee"):
        raw = str(payload.get(key) or "")
        match = re.search(r"(19|20)\d{2}", raw)
        if match:
            return match.group(0)
    return str(datetime.now().year)


# ── Rendu de secours (template Brand Box fond noir) ──────────────────────────

def _render_programmatic(payload: dict[str, Any]) -> str:
    prs, slide = _new_fiche_slide()
    g = _Grid(prs.slide_width, prs.slide_height)

    _add_bandeau(slide, g)
    _add_header_zone(slide, g, payload)
    _add_left_panel(slide, g, payload)
    _add_right_content(slide, g, payload)
    _add_footer(slide, g, payload)

    output_path = _build_output_path(payload.get("title"))
    prs.save(output_path)
    return str(output_path)


# ── Bandeau OCD ──────────────────────────────────────────────────────────────

def _add_bandeau(slide, g: _Grid) -> None:
    """Bandeau orange plein en haut (signature OCD)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, g.w, g.bandeau_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(g.margin, g.cm(0.1), g.cm(14), g.cm(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Orange Cyberdefense"
    p.font.name = _font()
    p.font.size = g.pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE


# ── Zone titre / sous-titre / badge ──────────────────────────────────────────

def _add_header_zone(slide, g: _Grid, payload: dict[str, Any]) -> None:
    """Zone de titre avec badge confidentialité, sur fond sombre du template."""
    y0 = g.bandeau_h
    title_text = _coerce(payload.get("title"), "Fiche Référence")
    confidentiality = payload.get("confidentiality", "Interne")

    # Titre principal — blanc sur fond noir
    txTitle = slide.shapes.add_textbox(
        g.margin, Emu(y0 + g.cm(0.35)), Emu(g.w - g.margin * 2 - g.cm(3.8)), g.cm(1.5)
    )
    tf = txTitle.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    s = _STYLES.get("title", {})
    p.font.name = _font(s.get("font_name"))
    p.font.size = g.pt(s.get("font_size", 26))
    p.font.bold = s.get("bold", True)
    p.font.color.rgb = TEXT_PRIMARY

    # Sous-titre : client | secteur | durée
    parts = [
        _coerce(payload.get("client"), ""),
        _coerce(payload.get("sector"), ""),
        _coerce(payload.get("duration"), ""),
    ]
    subtitle_text = "  |  ".join(p for p in parts if p)
    if not subtitle_text:
        subtitle_text = "[Client / Secteur / Durée à compléter]"

    txSub = slide.shapes.add_textbox(
        g.margin, Emu(y0 + g.cm(2.0)), Emu(g.w - g.margin * 2 - g.cm(3.8)), g.cm(0.9)
    )
    tf2 = txSub.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    s2 = _STYLES.get("subtitle", {})
    p2.font.name = _font(s2.get("font_name"))
    p2.font.size = g.pt(s2.get("font_size", 14))
    p2.font.bold = s2.get("bold", False)
    p2.font.color.rgb = TEXT_SECONDARY

    # Trait orange sous le titre
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, g.margin, Emu(y0 + g.cm(3.15)), Emu(g.w - g.margin * 2), g.cm(0.06)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = ORANGE
    sep.line.fill.background()

    _add_badge(slide, g, confidentiality)


def _add_badge(slide, g: _Grid, confidentiality: str) -> None:
    badge_w = g.cm(3.2)
    badge_h = g.cm(0.65)
    x = Emu(g.w - badge_w - g.margin)
    y = Emu(g.bandeau_h + g.cm(0.5))

    conf_hex = _CONF_COLORS.get(confidentiality, _COLORS.get("dark_grey", "#595959"))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, badge_w, badge_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(conf_hex)
    # Sur fond noir, un badge noir (Confidentiel) doit rester lisible : liseré orange.
    if _rgb(conf_hex) == DARK_BG:
        shape.line.color.rgb = ORANGE
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = confidentiality.upper()
    s = _STYLES.get("badge_text", {})
    p.font.name = _font(s.get("font_name"))
    p.font.size = g.pt(s.get("font_size", 9))
    p.font.bold = s.get("bold", True)
    p.font.color.rgb = WHITE


# ── Panneau gauche : métadonnées ─────────────────────────────────────────────

def _add_left_panel(slide, g: _Grid, payload: dict[str, Any]) -> None:
    """Colonne gauche sur fond gris sombre — infos structurées."""
    panel_x = g.margin
    panel_y = g.content_top
    panel_h = g.content_h

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, panel_x, panel_y, g.left_panel_w, panel_h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PANEL_FILL
    bg.line.color.rgb = PANEL_BORDER
    bg.line.width = Pt(0.5)

    rows = [
        ("CLIENT", _coerce(payload.get("client"), "[À COMPLÉTER]")),
        ("SECTEUR", _coerce(payload.get("sector"), "[À COMPLÉTER]")),
        ("DURÉE", _coerce(payload.get("duration"), "[À COMPLÉTER]")),
        ("ÉQUIPE", _coerce(payload.get("team"), "[À COMPLÉTER]")),
        ("MOTS-CLÉS", ", ".join(payload.get("keywords") or ["[À COMPLÉTER]"])),
    ]
    if payload.get("notes"):
        rows.append(("NOTE", _coerce(payload["notes"], "")))

    txBox = slide.shapes.add_textbox(
        Emu(panel_x + g.cm(0.35)),
        Emu(panel_y + g.cm(0.4)),
        Emu(g.left_panel_w - g.cm(0.7)),
        Emu(panel_h - g.cm(0.5)),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()

    label_s = _STYLES.get("metadata_label", {})
    value_s = _STYLES.get("metadata_value", {})

    for idx, (label, value) in enumerate(rows):
        if idx == 0:
            p_lbl = tf.paragraphs[0]
        else:
            sp = tf.add_paragraph()
            sp.text = ""
            sp.space_after = Pt(2)
            p_lbl = tf.add_paragraph()
        p_lbl.text = label
        p_lbl.font.name = _font(label_s.get("font_name"))
        p_lbl.font.size = g.pt(label_s.get("font_size", 9))
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = _rgb(label_s.get("color", "#FF7900"))
        p_lbl.space_after = Pt(1)

        p_val = tf.add_paragraph()
        p_val.text = value
        p_val.font.name = _font(value_s.get("font_name"))
        p_val.font.size = g.pt(value_s.get("font_size", 10))
        p_val.font.bold = False
        p_val.font.color.rgb = TEXT_PRIMARY
        p_val.space_after = Pt(4)


# ── Zone droite : sections de contenu ────────────────────────────────────────

def _add_right_content(slide, g: _Grid, payload: dict[str, Any]) -> None:
    """Zone principale à droite — Contexte, Mission, Livrables, Résultats."""
    x = g.right_x
    w = g.right_w
    y = g.content_top
    h = g.content_h

    half_h = Emu(int((h - g.cm(0.4)) / 2))
    col_w = Emu(int((w - g.cm(0.5)) / 2))

    _add_text_section(slide, g, "CONTEXTE", _coerce(payload.get("context"), "[À COMPLÉTER]"),
                      x, y, col_w, half_h)
    _add_text_section(slide, g, "MISSION", _coerce(payload.get("mission"), "[À COMPLÉTER]"),
                      Emu(x + col_w + g.cm(0.5)), y, col_w, half_h)

    sep_y = Emu(y + half_h + g.cm(0.15))
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, sep_y, w, g.cm(0.04))
    sep.fill.solid()
    sep.fill.fore_color.rgb = PANEL_BORDER
    sep.line.fill.background()

    y2 = Emu(sep_y + g.cm(0.2))
    h2 = Emu(g.h - y2 - g.cm(1.5))

    _add_bullet_section(slide, g, "LIVRABLES",
                        payload.get("deliverables") or ["[À COMPLÉTER]"],
                        x, y2, col_w, h2)
    _add_bullet_section(slide, g, "RÉSULTATS / VALEUR",
                        payload.get("results") or ["[À COMPLÉTER]"],
                        Emu(x + col_w + g.cm(0.5)), y2, col_w, h2)


def _add_text_section(slide, g: _Grid, heading: str, body: str, x, y, w, h) -> None:
    heading_s = _STYLES.get("heading", {})
    body_s = _STYLES.get("body", {})

    txH = slide.shapes.add_textbox(x, y, w, g.cm(0.6))
    p = txH.text_frame.paragraphs[0]
    p.text = heading
    p.font.name = _font(heading_s.get("font_name"))
    p.font.size = g.pt(heading_s.get("font_size", 11))
    p.font.bold = True
    p.font.color.rgb = _rgb(heading_s.get("color", "#FF7900"))

    txB = slide.shapes.add_textbox(x, Emu(y + g.cm(0.65)), w, Emu(h - g.cm(0.65)))
    tf = txB.text_frame
    tf.word_wrap = True
    p2 = tf.paragraphs[0]
    p2.text = body
    p2.font.name = _font(body_s.get("font_name"))
    p2.font.size = g.pt(body_s.get("font_size", 11))
    p2.font.bold = False
    p2.font.color.rgb = TEXT_PRIMARY


def _add_bullet_section(slide, g: _Grid, heading: str, items: list[str], x, y, w, h) -> None:
    heading_s = _STYLES.get("heading", {})
    bullet_s = _STYLES.get("bullet", {})

    txH = slide.shapes.add_textbox(x, y, w, g.cm(0.6))
    p = txH.text_frame.paragraphs[0]
    p.text = heading
    p.font.name = _font(heading_s.get("font_name"))
    p.font.size = g.pt(heading_s.get("font_size", 11))
    p.font.bold = True
    p.font.color.rgb = _rgb(heading_s.get("color", "#FF7900"))

    txB = slide.shapes.add_textbox(x, Emu(y + g.cm(0.65)), w, Emu(h - g.cm(0.65)))
    tf = txB.text_frame
    tf.word_wrap = True
    tf.clear()

    visible_items = (items or ["[À COMPLÉTER]"])[:3]
    for i, item in enumerate(visible_items):
        p2 = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p2.text = f"▸  {_coerce(item, '[À COMPLÉTER]')}"
        p2.font.name = _font(bullet_s.get("font_name"))
        p2.font.size = g.pt(bullet_s.get("font_size", 10))
        p2.font.bold = False
        p2.font.color.rgb = TEXT_PRIMARY
        p2.space_after = Pt(5)


# ── Footer ────────────────────────────────────────────────────────────────────

def _add_footer(slide, g: _Grid, payload: dict[str, Any]) -> None:
    footer_y = Emu(g.h - g.cm(1.3))

    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, g.margin, Emu(footer_y - g.cm(0.1)), Emu(g.w - g.margin * 2), g.cm(0.03)
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = PANEL_BORDER
    sep.line.fill.background()

    refs = payload.get("reference_examples_used") or []
    refs_text = f"Basé sur : {', '.join(refs[:3])}" if refs else ""
    date_str = datetime.now().strftime("%d/%m/%Y")

    left_text = "Orange Cyberdefense — Conseil & Audit"
    right_text = f"{refs_text}  |  REF-Factory  |  {date_str}" if refs_text else f"REF-Factory  |  {date_str}"

    footer_s = _STYLES.get("footer", {})

    txL = slide.shapes.add_textbox(g.margin, footer_y, g.cm(16), g.cm(0.7))
    pL = txL.text_frame.paragraphs[0]
    pL.text = left_text
    pL.font.name = _font(footer_s.get("font_name"))
    pL.font.size = g.pt(footer_s.get("font_size", 8))
    pL.font.color.rgb = MEDIUM_GREY

    txR = slide.shapes.add_textbox(Emu(g.w - g.cm(17) - g.margin), footer_y, g.cm(17), g.cm(0.7))
    pR = txR.text_frame.paragraphs[0]
    pR.text = right_text
    pR.alignment = PP_ALIGN.RIGHT
    pR.font.name = _font(footer_s.get("font_name"))
    pR.font.size = g.pt(footer_s.get("font_size", 8))
    pR.font.color.rgb = MEDIUM_GREY


# ── Utilitaires ──────────────────────────────────────────────────────────────

def _build_output_path(title: str | None) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    base = _slugify(title or "fiche-ref")
    return OUTPUT_DIR / f"{base}_{timestamp}.pptx"


def _slugify(value: str) -> str:
    normalized = unicodedata.normalize("NFKD", value)
    ascii_val = normalized.encode("ascii", "ignore").decode("ascii")
    ascii_val = re.sub(r"[^A-Za-z0-9]+", "-", ascii_val).strip("-")
    return ascii_val[:80] or "fiche-ref"


def _coerce(value: Any, default: str) -> str:
    if isinstance(value, str) and value.strip():
        return value.strip()
    return default
