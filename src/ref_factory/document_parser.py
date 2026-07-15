import json
import re
from pathlib import Path

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

SUPPORTED_TEXT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".json"}


def summarize_text(text: str, max_chars: int = 800) -> str:
    compact = re.sub(r"\s+", " ", text or "").strip()
    if len(compact) <= max_chars:
        return compact
    return f"{compact[: max_chars - 3].rstrip()}..."


def extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_text(file_path)
    if suffix == ".docx":
        return _extract_docx_text(file_path)
    if suffix == ".pptx":
        return _extract_pptx_text(file_path)
    if suffix in {".txt", ".md"}:
        return _extract_text_file(file_path)
    if suffix == ".json":
        return _extract_json_file(file_path)
    return ""


def parse_supporting_documents(file_paths: list[str]) -> list[dict[str, str]]:
    documents: list[dict[str, str]] = []
    for raw_path in file_paths or []:
        path = Path(raw_path)
        if not path.exists():
            continue
        text = extract_text_from_file(path)
        if not text.strip():
            continue
        documents.append(
            {
                "filename": path.name,
                "path": str(path),
                "extension": path.suffix.lower(),
                "text": text,
                "excerpt": summarize_text(text, max_chars=500),
            }
        )
    return documents


def build_combined_source_text(
    manual_fields: dict[str, str],
    brief_text: str,
    source_documents: list[dict[str, str]],
    max_chars: int = 12000,
) -> str:
    parts: list[str] = []

    manual_lines = [f"{key}: {value}" for key, value in (manual_fields or {}).items() if value]
    if manual_lines:
        parts.append("CHAMPS SAISIS\n" + "\n".join(manual_lines))

    if brief_text and brief_text.strip():
        parts.append("BRIEF LIBRE\n" + brief_text.strip())

    for document in source_documents:
        parts.append(f"DOCUMENT SOURCE - {document['filename']}\n{document['text']}")

    combined = "\n\n".join(parts).strip()
    if len(combined) <= max_chars:
        return combined
    return combined[:max_chars].rstrip()


def _extract_pdf_text(file_path: Path) -> str:
    try:
        reader = PdfReader(str(file_path))
    except Exception:
        return ""

    pages = []
    for page in reader.pages:
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        if text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx_text(file_path: Path) -> str:
    try:
        document = DocxDocument(str(file_path))
    except Exception:
        return ""

    blocks: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            blocks.append(text)

    for table in document.tables:
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_values:
                blocks.append(" | ".join(row_values))

    return "\n".join(blocks)


def _extract_pptx_text(file_path: Path) -> str:
    try:
        presentation = Presentation(str(file_path))
    except Exception:
        return ""

    slide_blocks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if isinstance(text, str) and text.strip():
                lines.append(text.strip())
        if len(lines) > 1:
            slide_blocks.append("\n".join(lines))
    return "\n\n".join(slide_blocks)


def _extract_text_file(file_path: Path) -> str:
    data = file_path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def _extract_json_file(file_path: Path) -> str:
    try:
        payload = json.loads(file_path.read_text(encoding="utf-8"))
    except Exception:
        return _extract_text_file(file_path)
    return json.dumps(payload, indent=2, ensure_ascii=True)
