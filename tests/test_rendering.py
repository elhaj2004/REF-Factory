from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor

from ref_factory.presentation.rendering import render_ref_slide


_PAYLOAD_COMPLET = {
    "title": "Audit de gouvernance SSI - Banque régionale",
    "client": "Banque X (anonymisé)",
    "sector": "Banque / Finance",
    "duration": "8 semaines",
    "team": "1 manager, 2 consultants",
    "keywords": ["audit", "gouvernance", "DORA", "ISO 27001"],
    "confidentiality": "Confidentiel",
    "context": "Évaluation de maturité SSI dans le cadre de la conformité DORA. Entité sans politique SSI formalisée.",
    "mission": "Audit de gouvernance SSI couvrant les 5 domaines DORA. Remise d'un rapport et d'une roadmap sur 18 mois.",
    "deliverables": ["Rapport audit gouvernance SSI", "Cartographie des écarts DORA", "Roadmap conformité 18 mois"],
    "results": ["Maturité 2.1/5 avant intervention", "32 recommandations priorisées", "Validation AMF du plan"],
    "notes": "Mission pilote dans le cadre du programme OCD Banque.",
    "reference_examples_used": ["REF_Banque_DORA.pptx", "REF_Audit_ISO27001.pptx"],
}

_PAYLOAD_MINIMAL = {
    "title": "Fiche REF minimale",
    "confidentiality": "Interne",
    "context": "Contexte test.",
    "mission": "Mission test.",
}


def test_render_creates_single_slide(tmp_path: Path) -> None:
    output_path = Path(render_ref_slide(_PAYLOAD_COMPLET))
    assert output_path.exists(), "Le fichier PPTX n'a pas été créé"
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1, "La présentation doit contenir exactement 1 slide"
    assert output_path.suffix.lower() == ".pptx"


def test_render_minimal_payload_does_not_crash() -> None:
    output_path = Path(render_ref_slide(_PAYLOAD_MINIMAL))
    assert output_path.exists()
    prs = Presentation(str(output_path))
    assert len(prs.slides) == 1


def test_slide_dimensions_template_fiche() -> None:
    """Vérifie les dimensions 33.87 × 19.05 cm du template fiche exact."""
    from pptx.util import Cm
    output_path = Path(render_ref_slide(_PAYLOAD_COMPLET))
    prs = Presentation(str(output_path))
    assert abs(prs.slide_width - Cm(33.87)) < 5000, "Largeur slide incorrecte"
    assert abs(prs.slide_height - Cm(19.05)) < 5000, "Hauteur slide incorrecte"


def test_render_uses_exact_fiche_template() -> None:
    """La fiche doit être un remplissage du template fiche exact : même layout,
    mêmes images (hors logo client retiré), quel que soit le secteur — seuls
    les textes changent."""
    from ref_factory.presentation.rendering import FICHE_TEMPLATE_PATH

    assert FICHE_TEMPLATE_PATH.exists(), f"Template fiche manquant : {FICHE_TEMPLATE_PATH}"
    template = Presentation(str(FICHE_TEMPLATE_PATH))
    template_pictures = sum(
        1 for sh in template.slides[0].shapes if sh.shape_type == 13
    )

    payload_banque = dict(_PAYLOAD_COMPLET)
    payload_sante = dict(_PAYLOAD_COMPLET, sector="Santé / Hôpital", title="PSSI Hôpital")

    for payload in (payload_banque, payload_sante):
        prs = Presentation(render_ref_slide(payload))
        slide = prs.slides[0]
        assert slide.slide_layout.name == template.slides[0].slide_layout.name, (
            "La fiche doit garder le layout du template fiche"
        )
        pictures = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        assert pictures == template_pictures - 1, (
            "Les images du template doivent être conservées (seul le logo client est retiré)"
        )
        all_text = " ".join(
            sh.text_frame.text for sh in slide.shapes if getattr(sh, "has_text_frame", False)
        )
        assert payload["title"] in all_text
        assert "Nexans" not in all_text, "Aucune trace du client d'origine ne doit rester"


def test_render_replaces_profile_pastilles() -> None:
    """Les pastilles du bloc 'Profil de la prestation' doivent refléter le payload."""
    prs = Presentation(render_ref_slide(_PAYLOAD_COMPLET))

    def walk(shapes):
        for sh in shapes:
            yield sh
            if sh.shape_type == 6:
                yield from walk(sh.shapes)

    all_text = " ".join(
        sh.text_frame.text for sh in walk(prs.slides[0].shapes)
        if getattr(sh, "has_text_frame", False)
    )
    assert "Banque / Finance" in all_text, "Le secteur doit apparaître dans la pastille"
    assert "8 semaines" in all_text, "La durée doit apparaître dans la pastille"
    assert "Banque X (anonymisé)" in all_text, "Le client doit remplacer le logo d'origine"


def test_slide_contains_title_text() -> None:
    """Le texte du titre doit apparaître quelque part dans la slide."""
    output_path = Path(render_ref_slide(_PAYLOAD_COMPLET))
    prs = Presentation(str(output_path))
    slide = prs.slides[0]
    all_text = " ".join(
        shape.text for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "Audit de gouvernance SSI" in all_text


def test_slide_keeps_template_section_headings() -> None:
    """Les intitulés de sections du template fiche doivent rester intacts."""
    output_path = Path(render_ref_slide(_PAYLOAD_COMPLET))
    prs = Presentation(str(output_path))
    slide = prs.slides[0]
    all_text = " ".join(
        shape.text for shape in slide.shapes if hasattr(shape, "text")
    )
    for heading in ("Contexte", "Réalisation", "Bénéfices"):
        assert heading in all_text, f"L'intitulé {heading!r} du template doit être conservé"
