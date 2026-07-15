"""Analyse la charte OCD depuis les templates Pres-Factory et les fiches REF existantes."""
import sys, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

sys.path.insert(0, str(Path(__file__).resolve().parent / "src"))

TEMPLATE_DIR = Path(r"C:\Users\BJPS1817\Pres-Factory\Tools and templates PPT - FR\French\2. Templates\French")
REF_SAMPLE = Path(r"O:\ConseiletAudit\Références\PC1 - Audit et Contrôle conformité\2025 - Cyberdiag Grand Est - PROCIVIS.pptx")

results = {}

def emu_to_inches(emu):
    return round(emu / 914400, 2) if emu else 0

def analyze_prs(prs, label):
    info = {}
    info["slide_width"] = emu_to_inches(prs.slide_width)
    info["slide_height"] = emu_to_inches(prs.slide_height)
    info["slide_count"] = len(prs.slides)
    
    # Analyser les masters
    masters_info = []
    for sm in prs.slide_masters:
        m = {}
        m["name"] = sm.name
        fonts = set()
        for layout in sm.slide_layouts:
            for ph in layout.placeholders:
                if hasattr(ph, 'text_frame'):
                    for p in ph.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
        m["fonts"] = list(fonts)
        masters_info.append(m)
    
    info["masters"] = masters_info
    
    # Analyser les slides
    slide_info = []
    for i, slide in enumerate(prs.slides[:3]):  # 3 premiers slides
        s = {"index": i+1}
        shapes = []
        for shape in slide.shapes:
            sh = {
                "type": str(shape.shape_type),
                "name": shape.name,
                "left": emu_to_inches(shape.left) if shape.left else 0,
                "top": emu_to_inches(shape.top) if shape.top else 0,
                "width": emu_to_inches(shape.width) if shape.width else 0,
                "height": emu_to_inches(shape.height) if shape.height else 0,
            }
            # Texte
            if hasattr(shape, 'text') and shape.text:
                sh["text"] = shape.text[:100]
            
            # Font et couleurs
            if hasattr(shape, 'text_frame'):
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.name:
                            sh.setdefault("font_name", run.font.name)
                        if run.font.size:
                            sh.setdefault("font_size_pt", run.font.size / 12700)
                        if run.font.color and run.font.color.rgb:
                            sh.setdefault("font_color", str(run.font.color.rgb))
                        if run.font.bold:
                            sh.setdefault("font_bold", True)
            
            # Remplissage shape
            if hasattr(shape, 'fill'):
                try:
                    if shape.fill.type is not None:
                        sh["fill_type"] = str(shape.fill.type)
                        try:
                            if shape.fill.fore_color and shape.fill.fore_color.rgb:
                                sh["fill_color"] = str(shape.fill.fore_color.rgb)
                        except:
                            pass
                except:
                    pass
            
            shapes.append(sh)
        
        s["shapes"] = shapes
        slide_info.append(s)
    
    info["slides_sample"] = slide_info
    return info

# 1. Analyser un template .potx
print("=== ANALYSE TEMPLATE OCD ===")
for potx in TEMPLATE_DIR.glob("*.potx"):
    print(f"\nAnalyse de {potx.name}...")
    try:
        prs = Presentation(str(potx))
        results[potx.name] = analyze_prs(prs, potx.name)
        print(f"  Masters: {len(prs.slide_masters)}, Slides: {len(prs.slides)}")
        print(f"  Dimensions: {results[potx.name]['slide_width']} x {results[potx.name]['slide_height']}")
        for sm in prs.slide_masters:
            print(f"  Master: {sm.name}")
            for layout in sm.slide_layouts[:5]:
                print(f"    Layout: {layout.name}")
                for ph in layout.placeholders:
                    print(f"      Placeholder {ph.placeholder_format.idx}: {ph.name} ({ph.placeholder_format.type})")
    except Exception as e:
        results[potx.name] = f"Error: {e}"
        print(f"  Error: {e}")

# 2. Analyser une fiche REF réelle
print(f"\n=== ANALYSE FICHE REF RÉELLE ===")
if REF_SAMPLE.exists():
    print(f"Analyse de {REF_SAMPLE.name}...")
    prs = Presentation(str(REF_SAMPLE))
    results["sample_ref"] = analyze_prs(prs, "sample_ref")
    print(f"  Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides[:5]):
        print(f"  Slide {i+1}:")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                print(f"    {shape.name}: '{shape.text[:80]}...'")
else:
    print(f"Fichier non trouvé: {REF_SAMPLE}")

# 3. Sauvegarder le résultat
output_path = Path(__file__).resolve().parent / "data" / "charter_analysis.json"
with open(output_path, 'w', encoding='utf-8') as f:
    json.dump(results, f, indent=2, ensure_ascii=False, default=str)

print(f"\nAnalyse sauvegardée dans {output_path}")
