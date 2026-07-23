"""
Supprime de data/reference_library/ les fiches REF (.pptx/.potx) qui ne
respectent pas la charte graphique Orange Cyberdefense.

Seuls les fichiers au format PowerPoint (.pptx, .potx) sont analysés : ce
sont les seuls a avoir des attributs visuels (dimensions, polices,
couleurs) verifiables vis-a-vis de la charte. Les autres formats
(.json, .txt, .md, .pdf, .docx) n'ont pas de "charte graphique" au sens
PowerPoint et ne sont pas concernes par ce script.

Criteres de conformite (calibres sur des fiches REF Orange CD reelles
validees comme exemples) :
  - Dimensions de slide : 25.40x14.29 cm OU 33.87x19.05 cm (+/- 0.05 cm),
    les deux formats etant utilises dans de vraies fiches REF.
  - Polices : Source Sans Pro, Calibri, Arial, Helvetica (toutes variantes),
    Montserrat. Les references de theme ("+mn-lt", "+mj-lt") sont ignorees
    (ce ne sont pas des noms de police, mais des pointeurs vers le theme).
  - Couleurs : nuances de gris (y compris noir/blanc), oranges de marque
    (#FF7900 officiel et #FF6600 legacy toleree), et l'accent bleu
    #00B0F0 observe dans des fiches reelles. Toute autre teinte (rouge,
    vert, violet, jaune vif...) indique un template hors charte OCD.

Un fichier est supprime des qu'un des criteres echoue clairement. Les
fichiers illisibles (ex. .ppt binaire non supporte par python-pptx) sont
signales mais jamais supprimes automatiquement, faute de pouvoir verifier
leur conformite.

Usage :
    python scripts/enforce_charter_compliance.py
"""
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from pptx import Presentation  # noqa: E402

DEST = Path(__file__).resolve().parents[1] / "data" / "reference_library"

CHECKABLE_EXTENSIONS = {".pptx", ".potx"}
UNREADABLE_EXTENSIONS = {".ppt"}  # binaire, non ouvrable par python-pptx

ALLOWED_DIMENSIONS_CM = [(25.40, 14.29), (33.87, 19.05)]
DIMENSION_TOLERANCE_CM = 0.05

ALLOWED_FONT_PREFIXES = ("source sans pro", "calibri", "arial", "helvetica", "montserrat")

ALLOWED_ACCENT_COLORS = {"FF7900", "FF6600", "00B0F0"}


def is_allowed_font(name: str) -> bool:
    if not name or name.startswith("+"):  # reference de theme, pas un nom de police
        return True
    lowered = name.strip().lower()
    return any(lowered.startswith(prefix) for prefix in ALLOWED_FONT_PREFIXES)


def is_allowed_color(hex_color: str) -> bool:
    hex_color = hex_color.upper()
    if hex_color in ALLOWED_ACCENT_COLORS:
        return True
    if len(hex_color) == 6:
        r, g, b = (int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        if abs(r - g) <= 5 and abs(g - b) <= 5 and abs(r - b) <= 5:
            return True  # nuance de gris (dont noir/blanc)
    return False


def matches_allowed_dimensions(width_cm: float, height_cm: float) -> bool:
    return any(
        abs(width_cm - w) <= DIMENSION_TOLERANCE_CM and abs(height_cm - h) <= DIMENSION_TOLERANCE_CM
        for w, h in ALLOWED_DIMENSIONS_CM
    )


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:  # groupe
            yield from walk_shapes(shape.shapes)


def collect_fonts_and_colors(prs: Presentation) -> tuple[set[str], set[str]]:
    fonts: set[str] = set()
    colors: set[str] = set()

    def add_color_from(obj) -> None:
        try:
            if obj.type is not None:
                colors.add(str(obj.fore_color.rgb))
        except Exception:
            pass

    for slide in prs.slides:
        for shape in walk_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)
                        try:
                            if run.font.color and run.font.color.type is not None:
                                colors.add(str(run.font.color.rgb))
                        except Exception:
                            pass
            try:
                add_color_from(shape.fill)
            except Exception:
                pass
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    fonts.add(run.font.name)
                        try:
                            add_color_from(cell.fill)
                        except Exception:
                            pass
    return fonts, colors


def check_compliance(file_path: Path) -> tuple[bool, list[str]]:
    """Retourne (conforme, liste des raisons de non-conformite)."""
    reasons: list[str] = []

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        return False, [f"fichier illisible ({e})"]

    width_cm = prs.slide_width / 360000
    height_cm = prs.slide_height / 360000
    if not matches_allowed_dimensions(width_cm, height_cm):
        reasons.append(f"dimensions non conformes ({width_cm:.2f} x {height_cm:.2f} cm)")

    fonts, colors = collect_fonts_and_colors(prs)

    bad_fonts = sorted({f for f in fonts if not is_allowed_font(f)})
    if bad_fonts:
        reasons.append(f"police(s) hors charte : {', '.join(bad_fonts)}")

    bad_colors = sorted({c for c in colors if not is_allowed_color(c)})
    if bad_colors:
        reasons.append(f"couleur(s) hors charte : {', '.join('#' + c for c in bad_colors)}")

    return (len(reasons) == 0), reasons


def main() -> int:
    if not DEST.exists():
        print(f"Dossier introuvable : {DEST}")
        return 1

    checked = 0
    removed = 0
    skipped_unreadable = 0
    skipped_out_of_scope = 0

    for file_path in sorted(DEST.rglob("*")):
        if not file_path.is_file() or file_path.name.startswith("."):
            continue

        suffix = file_path.suffix.lower()

        if suffix in UNREADABLE_EXTENSIONS:
            print(f"[VERIF MANUELLE] {file_path.name} : format .ppt non analysable, non supprimé automatiquement.")
            skipped_unreadable += 1
            continue

        if suffix not in CHECKABLE_EXTENSIONS:
            skipped_out_of_scope += 1
            continue

        checked += 1
        compliant, reasons = check_compliance(file_path)
        if compliant:
            print(f"[OK] {file_path.name}")
        else:
            print(f"[SUPPRIME] {file_path.name} : {' | '.join(reasons)}")
            file_path.unlink()
            removed += 1

    print(
        f"\nTerminé : {checked} fichier(s) PowerPoint analysé(s), {removed} supprimé(s), "
        f"{skipped_unreadable} illisible(s) (à vérifier manuellement), "
        f"{skipped_out_of_scope} hors périmètre (formats non PowerPoint, non concernés par la charte graphique)."
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
